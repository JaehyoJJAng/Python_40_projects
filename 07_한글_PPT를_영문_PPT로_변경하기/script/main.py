import googletrans
from pptx import Presentation
from pptx.util import Pt
import os
import glob

class Translate:
    def __init__(self) -> None:
        # GoogleTrans Instance
        self.google_translator = googletrans.Translator()
    
    def translate(self,text:str) -> str:
        """ text 매개변수를 영어로 번역 """
        return str(self.google_translator.translate(text=text,dest='en',src='auto').text)

class PPTX:
    def __init__(self,filePath:str,translate)-> None:
        self.filePath = filePath
        self.trans = translate

    def trans_making_new_pptx(self)-> None:
        """ PPTX 파일 읽어온 후 영문번역 - 파일 저장 """
        # Presentation Instance
        prs = Presentation(self.filePath)

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == 14 or shape.shape_type == 17:
                    text = shape.text

                    # 번역
                    translate_text : str = self.trans.translate(text=text)

                    # shape를 번역된 내용으로 변환
                    shape.text_frame.text = translate_text

        # 저장 경로 지정
        savePath : str = '../result'
        fileName : str = os.path.join(savePath,'영문번환.pptx')
        if not os.path.exists(savePath):
            os.mkdir(savePath)

        # 저장하기
        prs.save(fileName)

        print(f'{fileName} - 저장 완료!')

def main()-> None:
    # pptx File Path
    pptx_path = glob.glob('../pptx/*.pptx')[0]

    # Google Translate Instance
    trans = Translate()

    # PPTX
    pptx = PPTX(filePath=pptx_path,translate=trans)

    # 번역 및 저장
    pptx.trans_making_new_pptx()

if __name__ == '__main__':
    main() 